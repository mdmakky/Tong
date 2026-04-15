"""
Generate a comprehensive PPTX documentation for the Tong Messaging Platform.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Constants ───────────────────────────────────────────────────────────
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Colors
BG_DARK = RGBColor(0x0F, 0x10, 0x17)
BG_CARD = RGBColor(0x1A, 0x1B, 0x26)
BG_CARD2 = RGBColor(0x23, 0x24, 0x33)
ACCENT_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
ACCENT_INDIGO = RGBColor(0x63, 0x66, 0xF1)
ACCENT_BLUE = RGBColor(0x38, 0xBD, 0xF8)
ACCENT_GREEN = RGBColor(0x22, 0xC5, 0x5E)
ACCENT_AMBER = RGBColor(0xF5, 0x9E, 0x0B)
ACCENT_RED = RGBColor(0xEF, 0x44, 0x44)
ACCENT_PINK = RGBColor(0xEC, 0x48, 0x99)
TEXT_WHITE = RGBColor(0xF0, 0xF0, 0xF5)
TEXT_GRAY = RGBColor(0x94, 0xA3, 0xB8)
TEXT_LIGHT_GRAY = RGBColor(0xCB, 0xD5, 0xE1)
BORDER_COLOR = RGBColor(0x33, 0x34, 0x45)


def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or BG_CARD
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    # Rounded corners
    shape.adjustments[0] = 0.05
    return shape


def add_rect(slide, left, top, width, height, fill_color=None, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or BG_CARD
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_circle(slide, left, top, size, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14, color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_text(slide, left, top, width, height, lines, font_size=13, color=TEXT_WHITE, line_spacing=1.5, font_name="Segoe UI"):
    """lines is a list of (text, bold, color) tuples"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, bold, col) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = col or color
        p.font.bold = bold
        p.font.name = font_name
        p.space_after = Pt(4)
    return txBox


def add_arrow(slide, left, top, width, height, fill_color=ACCENT_PURPLE):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_right_arrow(slide, left, top, width, height, fill_color=ACCENT_PURPLE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_chevron(slide, left, top, width, height, fill_color=ACCENT_PURPLE):
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


# ─── Slide Builders ──────────────────────────────────────────────────────

def create_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide)

    # Decorative accent bar at top
    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ACCENT_PURPLE)

    # Decorative circles
    add_circle(slide, Inches(10.5), Inches(1.2), Inches(2.5), RGBColor(0x8B, 0x5C, 0xF6))
    add_circle(slide, Inches(11.2), Inches(2.0), Inches(1.8), RGBColor(0x63, 0x66, 0xF1))

    # Title
    add_text_box(slide, Inches(1), Inches(1.5), Inches(9), Inches(1.2),
                 "💬 tong", font_size=60, color=TEXT_WHITE, bold=True)

    add_text_box(slide, Inches(1), Inches(2.8), Inches(9), Inches(0.8),
                 "Real-time Messaging Platform", font_size=32, color=ACCENT_BLUE, bold=True)

    add_text_box(slide, Inches(1), Inches(3.8), Inches(9), Inches(0.6),
                 "Complete Project Documentation — Phase 1: Messaging System",
                 font_size=18, color=TEXT_GRAY)

    # Info cards
    info_items = [
        ("🛠️  Tech Stack", "React.js + Node.js / Express.js"),
        ("📦  Databases", "PostgreSQL + MongoDB + Redis"),
        ("⚡  Real-time", "Socket.io WebSocket"),
        ("🚀  Deployment", "Docker / Render / Vercel"),
    ]
    for i, (label, value) in enumerate(info_items):
        x = Inches(1) + Inches(i * 2.8)
        card = add_shape(slide, x, Inches(5.0), Inches(2.5), Inches(1.5), BG_CARD, BORDER_COLOR)
        add_text_box(slide, x + Inches(0.2), Inches(5.15), Inches(2.1), Inches(0.5),
                     label, font_size=12, color=ACCENT_PURPLE, bold=True)
        add_text_box(slide, x + Inches(0.2), Inches(5.6), Inches(2.1), Inches(0.7),
                     value, font_size=13, color=TEXT_WHITE)


def create_toc_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ACCENT_INDIGO)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                 "📋 Table of Contents", font_size=32, color=TEXT_WHITE, bold=True)

    toc_items = [
        "1.  Project Overview & 5-Phase Roadmap",
        "2.  System Architecture Overview",
        "3.  Tech Stack — Frontend & Backend",
        "4.  Hybrid Database Architecture",
        "5.  Database Schema Design (PostgreSQL + MongoDB + Redis)",
        "6.  Authentication & Security Workflow",
        "7.  Real-time Messaging — Socket.io Flow",
        "8.  One-to-One Chat Feature Workflow",
        "9.  Group Chat Feature Workflow",
        "10. File & Media Upload Workflow",
        "11. Message Lifecycle (Edit, Delete, Pin, React)",
        "12. Notification & Presence System",
        "13. Frontend Architecture & Component Tree",
        "14. API Endpoint Summary",
        "15. Deployment — Server Architecture (Live Droplet)",
        "16. Deployment — Deploy Flow & External Services",
        "17. Project File Structure",
    ]

    col1 = toc_items[:9]
    col2 = toc_items[9:]

    for i, item in enumerate(col1):
        add_text_box(slide, Inches(1.2), Inches(1.4 + i * 0.6), Inches(5.5), Inches(0.5),
                     item, font_size=15, color=TEXT_LIGHT_GRAY)

    for i, item in enumerate(col2):
        add_text_box(slide, Inches(7.2), Inches(1.4 + i * 0.6), Inches(5.5), Inches(0.5),
                     item, font_size=15, color=TEXT_LIGHT_GRAY)


def create_overview_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ACCENT_GREEN)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                 "1. Project Overview", font_size=32, color=TEXT_WHITE, bold=True)

    # Description
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(1.0),
                 "tong is a full-featured, real-time messaging platform built with React.js & Node.js/Express.js. "
                 "The project is planned in 5 phases, with Phase 1 delivering a complete messaging system that serves as the foundation for all future phases.",
                 font_size=14, color=TEXT_GRAY)

    # 5-Phase Roadmap
    add_text_box(slide, Inches(0.8), Inches(2.3), Inches(10), Inches(0.5),
                 "5-Phase Roadmap", font_size=20, color=ACCENT_BLUE, bold=True)

    phases = [
        ("Phase 1", "Complete Messaging\n(1-to-1, Group, Private)", ACCENT_GREEN, "🟢 CURRENT"),
        ("Phase 2", "Audio Calling\n(WebRTC-based)", ACCENT_BLUE, "⏳ Planned"),
        ("Phase 3", "Video Calling\n(WebRTC + Media Server)", ACCENT_PURPLE, "⏳ Planned"),
        ("Phase 4", "Anonymous Chat\n(Omegle-style)", ACCENT_AMBER, "⏳ Planned"),
        ("Phase 5", "Mobile App\n(React Native)", ACCENT_PINK, "⏳ Planned"),
    ]

    for i, (name, desc, color, status) in enumerate(phases):
        x = Inches(0.8) + Inches(i * 2.45)
        card = add_shape(slide, x, Inches(3.0), Inches(2.2), Inches(2.4), BG_CARD, color)
        add_text_box(slide, x + Inches(0.15), Inches(3.15), Inches(1.9), Inches(0.4),
                     name, font_size=16, color=color, bold=True)
        add_text_box(slide, x + Inches(0.15), Inches(3.6), Inches(1.9), Inches(1.0),
                     desc, font_size=12, color=TEXT_LIGHT_GRAY)
        add_text_box(slide, x + Inches(0.15), Inches(4.7), Inches(1.9), Inches(0.4),
                     status, font_size=11, color=color, bold=True)

        # Arrow between phases
        if i < len(phases) - 1:
            add_right_arrow(slide, x + Inches(2.2), Inches(3.9), Inches(0.25), Inches(0.3), color)

    # Phase 1 Goals
    add_text_box(slide, Inches(0.8), Inches(5.7), Inches(10), Inches(0.4),
                 "Phase 1 Goals:", font_size=16, color=ACCENT_GREEN, bold=True)

    goals = [
        "✅ Real-time One-to-One & Group Chat",
        "✅ JWT Auth + 2FA + OTP Email Verification",
        "✅ File/Media Sharing via Cloudinary",
        "✅ Message Edit, Delete, Reply, Forward, Pin, Reactions",
        "✅ Typing Indicators & Read Receipts",
        "✅ Scalable Architecture for future phases",
    ]
    for i, goal in enumerate(goals):
        col = i // 3
        row = i % 3
        add_text_box(slide, Inches(0.8 + col * 5.5), Inches(6.2 + row * 0.35), Inches(5.3), Inches(0.35),
                     goal, font_size=12, color=TEXT_LIGHT_GRAY)


def create_architecture_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ACCENT_BLUE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                 "2. System Architecture Overview", font_size=32, color=TEXT_WHITE, bold=True)

    # Client Layer
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(3), Inches(0.4),
                 "CLIENT LAYER", font_size=14, color=ACCENT_BLUE, bold=True)
    card = add_shape(slide, Inches(0.8), Inches(1.7), Inches(3.5), Inches(1.6), BG_CARD, ACCENT_BLUE)
    add_text_box(slide, Inches(1.0), Inches(1.85), Inches(3.0), Inches(0.35),
                 "🖥️  React.js 18+ (Vite)", font_size=14, color=TEXT_WHITE, bold=True)
    items = ["Zustand State Management", "Socket.io-client (WebSocket)", "TanStack Query (API Cache)", "React Router v6 (SPA)"]
    for i, item in enumerate(items):
        add_text_box(slide, Inches(1.2), Inches(2.25 + i * 0.25), Inches(2.8), Inches(0.25),
                     f"• {item}", font_size=10, color=TEXT_GRAY)

    # Arrows
    add_right_arrow(slide, Inches(4.5), Inches(2.2), Inches(0.5), Inches(0.3), ACCENT_BLUE)
    add_right_arrow(slide, Inches(4.5), Inches(2.7), Inches(0.5), Inches(0.3), ACCENT_GREEN)

    # Labels
    add_text_box(slide, Inches(4.3), Inches(1.85), Inches(0.9), Inches(0.3),
                 "REST", font_size=9, color=ACCENT_BLUE, bold=True)
    add_text_box(slide, Inches(4.2), Inches(2.95), Inches(1.0), Inches(0.3),
                 "WebSocket", font_size=9, color=ACCENT_GREEN, bold=True)

    # Server Layer
    add_text_box(slide, Inches(5.3), Inches(1.3), Inches(3), Inches(0.4),
                 "SERVER LAYER", font_size=14, color=ACCENT_GREEN, bold=True)
    card = add_shape(slide, Inches(5.3), Inches(1.7), Inches(3.5), Inches(1.6), BG_CARD, ACCENT_GREEN)
    add_text_box(slide, Inches(5.5), Inches(1.85), Inches(3.0), Inches(0.35),
                 "⚙️  Node.js + Express.js", font_size=14, color=TEXT_WHITE, bold=True)
    items = ["REST API (5 route modules)", "Socket.io Server (3 handlers)", "JWT Authentication Middleware", "Rate Limiting + Helmet Security"]
    for i, item in enumerate(items):
        add_text_box(slide, Inches(5.7), Inches(2.25 + i * 0.25), Inches(2.8), Inches(0.25),
                     f"• {item}", font_size=10, color=TEXT_GRAY)

    # Arrows to DBs
    add_arrow(slide, Inches(5.8), Inches(3.5), Inches(0.3), Inches(0.4), ACCENT_PURPLE)
    add_arrow(slide, Inches(6.8), Inches(3.5), Inches(0.3), Inches(0.4), ACCENT_GREEN)
    add_arrow(slide, Inches(7.8), Inches(3.5), Inches(0.3), Inches(0.4), ACCENT_RED)

    # Data Layer
    add_text_box(slide, Inches(0.8), Inches(3.7), Inches(3), Inches(0.4),
                 "DATA LAYER", font_size=14, color=ACCENT_PURPLE, bold=True)

    # PostgreSQL
    pg = add_shape(slide, Inches(0.8), Inches(4.2), Inches(3.5), Inches(1.5), BG_CARD, ACCENT_PURPLE)
    add_text_box(slide, Inches(1.0), Inches(4.35), Inches(3.0), Inches(0.35),
                 "🐘 PostgreSQL (Prisma ORM)", font_size=13, color=ACCENT_PURPLE, bold=True)
    pg_items = ["Users, Groups, Conversations", "Friend Requests, Contacts", "User Devices, Blocks, Reports", "ACID compliance + Complex joins"]
    for i, item in enumerate(pg_items):
        add_text_box(slide, Inches(1.2), Inches(4.7 + i * 0.24), Inches(2.8), Inches(0.24),
                     f"• {item}", font_size=10, color=TEXT_GRAY)

    # MongoDB
    mg = add_shape(slide, Inches(4.8), Inches(4.2), Inches(3.5), Inches(1.5), BG_CARD, ACCENT_GREEN)
    add_text_box(slide, Inches(5.0), Inches(4.35), Inches(3.0), Inches(0.35),
                 "🍃 MongoDB (Mongoose ODM)", font_size=13, color=ACCENT_GREEN, bold=True)
    mg_items = ["Messages collection (all types)", "Reactions, Read Receipts", "Edit History, Delivery Status", "TTL Index for disappearing msgs"]
    for i, item in enumerate(mg_items):
        add_text_box(slide, Inches(5.2), Inches(4.7 + i * 0.24), Inches(2.8), Inches(0.24),
                     f"• {item}", font_size=10, color=TEXT_GRAY)

    # Redis
    rd = add_shape(slide, Inches(8.8), Inches(4.2), Inches(3.5), Inches(1.5), BG_CARD, ACCENT_RED)
    add_text_box(slide, Inches(9.0), Inches(4.35), Inches(3.0), Inches(0.35),
                 "⚡ Redis (ioredis)", font_size=13, color=ACCENT_RED, bold=True)
    rd_items = ["OTP storage (10 min TTL)", "Token blacklist (on logout)", "Session cache", "In-memory fallback available"]
    for i, item in enumerate(rd_items):
        add_text_box(slide, Inches(9.2), Inches(4.7 + i * 0.24), Inches(2.8), Inches(0.24),
                     f"• {item}", font_size=10, color=TEXT_GRAY)

    # External Services
    add_text_box(slide, Inches(9.5), Inches(1.3), Inches(3.5), Inches(0.4),
                 "EXTERNAL SERVICES", font_size=14, color=ACCENT_AMBER, bold=True)
    ext = add_shape(slide, Inches(9.5), Inches(1.7), Inches(3.5), Inches(1.6), BG_CARD, ACCENT_AMBER)
    ext_items = [
        ("☁️  Cloudinary", "File/Image/Video Storage"),
        ("📧  Brevo API", "Transactional Email (OTP)"),
        ("🔐  JWT + bcrypt", "Manual Auth (no 3rd party)"),
        ("📱  QR Code (qrcode)", "Group Invite QR"),
    ]
    for i, (label, desc) in enumerate(ext_items):
        add_text_box(slide, Inches(9.7), Inches(1.85 + i * 0.35), Inches(3.0), Inches(0.35),
                     f"{label}: {desc}", font_size=10, color=TEXT_GRAY)

    # Data flow summary
    add_shape(slide, Inches(0.8), Inches(6.1), Inches(11.7), Inches(1.1), BG_CARD2, BORDER_COLOR)
    add_text_box(slide, Inches(1.0), Inches(6.2), Inches(11.3), Inches(0.3),
                 "🔄 Data Flow Summary", font_size=14, color=ACCENT_BLUE, bold=True)
    add_text_box(slide, Inches(1.0), Inches(6.55), Inches(11.3), Inches(0.5),
                 "React Client ──(REST API)──▶ Express.js ──▶ PostgreSQL (users/groups) & MongoDB (messages) & Redis (cache/OTP)  |  "
                 "React Client ──(WebSocket)──▶ Socket.io Server ──▶ Real-time events (new_message, typing, presence, reactions)",
                 font_size=11, color=TEXT_GRAY)


def create_tech_stack_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ACCENT_PURPLE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                 "3. Tech Stack — Frontend & Backend", font_size=32, color=TEXT_WHITE, bold=True)

    # Frontend
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.4),
                 "Frontend Stack", font_size=20, color=ACCENT_BLUE, bold=True)

    fe_items = [
        ("React.js 18+", "UI Framework — Component-based, fast rendering"),
        ("Vite", "Build tool — lightning-fast HMR dev server"),
        ("Zustand", "State management — lightweight, hook-based store"),
        ("Socket.io-client", "WebSocket — real-time bidirectional events"),
        ("TanStack Query", "Server state — API caching & sync"),
        ("React Router v6", "SPA routing — protected & public routes"),
        ("Axios", "HTTP client — REST API calls with interceptors"),
        ("TailwindCSS 3", "Utility-first CSS framework"),
        ("Lucide React", "Icon library — clean SVG icons"),
        ("Emoji Picker", "Emoji selection for message reactions"),
        ("React Dropzone", "Drag-and-drop file upload"),
        ("React Virtuoso", "Virtual lists for message performance"),
        ("React Hot Toast", "Toast notifications"),
        ("Headless UI", "Accessible UI primitives (modals, menus)"),
    ]

    for i, (name, desc) in enumerate(fe_items):
        col = i // 7
        row = i % 7
        x = Inches(0.8 + col * 6)
        y = Inches(1.8 + row * 0.7)
        add_shape(slide, x, y, Inches(5.5), Inches(0.6), BG_CARD, BORDER_COLOR)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.05), Inches(2.0), Inches(0.25),
                     name, font_size=12, color=ACCENT_BLUE, bold=True)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.3), Inches(5.1), Inches(0.25),
                     desc, font_size=10, color=TEXT_GRAY)

    # Backend
    add_text_box(slide, Inches(6.8), Inches(1.3), Inches(5.5), Inches(0.4),
                 "Backend Stack", font_size=20, color=ACCENT_GREEN, bold=True)

    be_items = [
        ("Node.js 18+", "Runtime — Non-blocking I/O, ES Modules"),
        ("Express.js", "HTTP Framework — REST API with routing"),
        ("Socket.io", "WebSocket server — real-time messaging"),
        ("Prisma ORM", "PostgreSQL ORM — type-safe queries, migrations"),
        ("Mongoose", "MongoDB ODM — schema-based message models"),
        ("ioredis", "Redis client — OTP, token blacklisting"),
        ("JWT + bcrypt", "Auth — stateless tokens, password hashing"),
        ("Multer + Cloudinary", "File upload — middleware + cloud storage"),
        ("Brevo API", "Email — transactional OTP emails"),
        ("Speakeasy + QRCode", "2FA — TOTP generation & verification"),
        ("express-validator", "Input validation & sanitization"),
        ("Helmet + CORS", "Security headers & CORS policy"),
        ("express-rate-limit", "Rate limiting — brute force protection"),
        ("Morgan", "HTTP request logging"),
    ]

    for i, (name, desc) in enumerate(be_items):
        col = i // 7
        row = i % 7
        x = Inches(6.8) if col == 0 else Inches(0.8)
        y = Inches(1.8 + row * 0.7) if col == 0 else Inches(1.8 + row * 0.7)
        if col == 1:
            continue  # We already have 14 items in 2 columns
        add_shape(slide, x, y, Inches(5.5), Inches(0.6), BG_CARD, BORDER_COLOR)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.05), Inches(2.0), Inches(0.25),
                     name, font_size=12, color=ACCENT_GREEN, bold=True)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.3), Inches(5.1), Inches(0.25),
                     desc, font_size=10, color=TEXT_GRAY)


def create_db_schema_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ACCENT_PURPLE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                 "4–5. Database Schema Design", font_size=32, color=TEXT_WHITE, bold=True)

    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.4),
                 "Hybrid Database Strategy: PostgreSQL for relational data, MongoDB for messages, Redis for cache",
                 font_size=13, color=TEXT_GRAY)

    # PostgreSQL Tables
    add_text_box(slide, Inches(0.8), Inches(1.6), Inches(4), Inches(0.4),
                 "🐘 PostgreSQL (Prisma) — 8 Tables", font_size=16, color=ACCENT_PURPLE, bold=True)

    pg_tables = [
        ("User", "id, username, email, password_hash,\ndisplay_name, avatar_url, bio, status,\nis_verified, 2FA, theme, language"),
        ("Group", "id, name, unique_group_id, type,\nowner_id, max_members, invite_link,\nis_invite_only, message_retention"),
        ("GroupMember", "group_id, user_id, role (owner/\nadmin/moderator/member),\nnickname, muted_until, last_read_at"),
        ("Conversation", "id, type (direct/private_encrypted),\nparticipant_1, participant_2,\nis_blocked, encryption_key_hash"),
        ("UserBlock", "blocker_id, blocked_id\n(unique constraint pair)"),
        ("UserReport", "reporter_id, reported_id,\nreason, status (pending/reviewed)"),
        ("UserDevice", "user_id, device_name, device_type,\nip_address, refresh_token, last_active"),
        ("FriendRequest", "sender_id, receiver_id,\nstatus (pending/accepted/rejected)"),
    ]

    for i, (name, fields) in enumerate(pg_tables):
        col = i // 4
        row = i % 4
        x = Inches(0.8 + col * 3.2)
        y = Inches(2.1 + row * 1.25)
        card = add_shape(slide, x, y, Inches(2.9), Inches(1.15), BG_CARD, ACCENT_PURPLE)
        add_text_box(slide, x + Inches(0.1), y + Inches(0.05), Inches(2.7), Inches(0.25),
                     name, font_size=12, color=ACCENT_PURPLE, bold=True)
        add_text_box(slide, x + Inches(0.1), y + Inches(0.3), Inches(2.7), Inches(0.8),
                     fields, font_size=9, color=TEXT_GRAY)

    # MongoDB Collection
    add_text_box(slide, Inches(7.3), Inches(1.6), Inches(5.5), Inches(0.4),
                 "🍃 MongoDB — Messages Collection", font_size=16, color=ACCENT_GREEN, bold=True)

    mongo_card = add_shape(slide, Inches(7.3), Inches(2.1), Inches(5.5), Inches(3.5), BG_CARD, ACCENT_GREEN)
    mongo_fields = [
        "conversation_id → links to PG conversation.id or group.id",
        "conversation_type → 'direct' | 'group' | 'private'",
        "sender_id → links to PG user.id",
        "message_type → text | image | video | audio | file | sticker | location | system",
        "content { text, media_url, media_type, media_size, thumbnail_url,",
        "          file_name, duration, location, is_encrypted }",
        "reply_to → ObjectId ref to another Message",
        "reactions [ { user_id, emoji, reacted_at } ]",
        "read_receipts [ { user_id, read_at } ]",
        "delivered_to [ user_ids ]",
        "is_pinned, is_edited, edit_history [ { content, edited_at } ]",
        "deleted_for [ user_ids ], deleted_for_all",
        "expires_at → TTL index for disappearing messages",
        "mentions [ user_ids ]",
    ]
    for i, field in enumerate(mongo_fields):
        add_text_box(slide, Inches(7.5), Inches(2.25 + i * 0.23), Inches(5.1), Inches(0.23),
                     field, font_size=9, color=TEXT_GRAY, font_name="Consolas")

    # Redis
    add_text_box(slide, Inches(7.3), Inches(5.8), Inches(5.5), Inches(0.4),
                 "⚡ Redis — Key-Value Cache", font_size=16, color=ACCENT_RED, bold=True)

    redis_card = add_shape(slide, Inches(7.3), Inches(6.2), Inches(5.5), Inches(1.0), BG_CARD, ACCENT_RED)
    redis_items = [
        "otp:{type}:{email}  →  6-digit code  (TTL: 10 min)",
        "bl:{token}          →  '1'           (TTL: remaining token exp)",
        "In-memory Map fallback when Redis is unavailable",
    ]
    for i, item in enumerate(redis_items):
        add_text_box(slide, Inches(7.5), Inches(6.3 + i * 0.22), Inches(5.1), Inches(0.22),
                     item, font_size=10, color=TEXT_GRAY, font_name="Consolas")


def create_auth_workflow_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ACCENT_AMBER)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                 "6. Authentication & Security Workflow", font_size=32, color=TEXT_WHITE, bold=True)

    # Registration Flow
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.4),
                 "📝 Registration Flow", font_size=18, color=ACCENT_GREEN, bold=True)

    reg_steps = [
        ("1", "User submits\nusername, email, password", ACCENT_BLUE),
        ("2", "Password hashed\n(bcrypt, cost=12)", ACCENT_PURPLE),
        ("3", "User record created\nin PostgreSQL", ACCENT_GREEN),
        ("4", "OTP generated &\nstored in Redis (10m)", ACCENT_AMBER),
        ("5", "OTP email sent\nvia Brevo API", ACCENT_PINK),
        ("6", "JWT tokens returned\n(access + refresh)", ACCENT_BLUE),
    ]

    for i, (num, text, color) in enumerate(reg_steps):
        x = Inches(0.8 + i * 2.0)
        card = add_shape(slide, x, Inches(1.8), Inches(1.75), Inches(1.3), BG_CARD, color)
        circle = add_circle(slide, x + Inches(0.6), Inches(1.65), Inches(0.35), color)
        add_text_box(slide, x + Inches(0.7), Inches(1.7), Inches(0.2), Inches(0.3),
                     num, font_size=12, color=BG_DARK, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.1), Inches(2.1), Inches(1.55), Inches(0.9),
                     text, font_size=10, color=TEXT_LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

        if i < len(reg_steps) - 1:
            add_right_arrow(slide, x + Inches(1.75), Inches(2.3), Inches(0.25), Inches(0.2), color)

    # Login Flow
    add_text_box(slide, Inches(0.8), Inches(3.4), Inches(5.5), Inches(0.4),
                 "🔑 Login Flow", font_size=18, color=ACCENT_BLUE, bold=True)

    login_steps = [
        ("1", "User submits\nemail + password", ACCENT_BLUE),
        ("2", "Validate credentials\n(bcrypt.compare)", ACCENT_PURPLE),
        ("3", "Check 2FA status\n(if enabled → TOTP)", ACCENT_AMBER),
        ("4", "Generate token pair\nAccess(15m) + Refresh(7d)", ACCENT_GREEN),
        ("5", "Store device info\n(IP, UA, refresh token)", ACCENT_PINK),
        ("6", "Initialize Socket.io\nconnection on client", ACCENT_BLUE),
    ]

    for i, (num, text, color) in enumerate(login_steps):
        x = Inches(0.8 + i * 2.0)
        card = add_shape(slide, x, Inches(3.9), Inches(1.75), Inches(1.3), BG_CARD, color)
        circle = add_circle(slide, x + Inches(0.6), Inches(3.75), Inches(0.35), color)
        add_text_box(slide, x + Inches(0.7), Inches(3.8), Inches(0.2), Inches(0.3),
                     num, font_size=12, color=BG_DARK, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.1), Inches(4.2), Inches(1.55), Inches(0.9),
                     text, font_size=10, color=TEXT_LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

        if i < len(login_steps) - 1:
            add_right_arrow(slide, x + Inches(1.75), Inches(4.4), Inches(0.25), Inches(0.2), color)

    # Token Management
    add_text_box(slide, Inches(0.8), Inches(5.5), Inches(5), Inches(0.4),
                 "🔐 Token & Security Details", font_size=18, color=ACCENT_RED, bold=True)

    sec_card = add_shape(slide, Inches(0.8), Inches(5.9), Inches(5.5), Inches(1.3), BG_CARD, BORDER_COLOR)
    sec_items = [
        "• Access Token: JWT, 15 min expiry, sent in Authorization header",
        "• Refresh Token: JWT, 7 day expiry, stored per-device",
        "• Token Rotation: New refresh token on each refresh",
        "• Blacklisting: Logout invalidates tokens in Redis",
        "• Password: bcrypt hash with cost factor 12",
    ]
    for i, item in enumerate(sec_items):
        add_text_box(slide, Inches(1.0), Inches(6.0 + i * 0.22), Inches(5.1), Inches(0.22),
                     item, font_size=10, color=TEXT_GRAY)

    # 2FA Flow
    sec2_card = add_shape(slide, Inches(6.8), Inches(5.9), Inches(5.5), Inches(1.3), BG_CARD, BORDER_COLOR)
    add_text_box(slide, Inches(7.0), Inches(5.95), Inches(5.0), Inches(0.3),
                 "🔒 Two-Factor Authentication (2FA)", font_size=13, color=ACCENT_AMBER, bold=True)
    twofa_items = [
        "• TOTP via Speakeasy (compatible with Google Authenticator)",
        "• Enable: Server generates secret → QR code → User scans",
        "• Verify: User submits 6-digit TOTP code to confirm setup",
        "• Login: If 2FA enabled, TOTP code required after password",
    ]
    for i, item in enumerate(twofa_items):
        add_text_box(slide, Inches(7.0), Inches(6.3 + i * 0.22), Inches(5.1), Inches(0.22),
                     item, font_size=10, color=TEXT_GRAY)


def create_socket_flow_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ACCENT_GREEN)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                 "7. Real-time Messaging — Socket.io Workflow", font_size=32, color=TEXT_WHITE, bold=True)

    # Connection Flow
    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(12), Inches(0.4),
                 "⚡ Connection Lifecycle", font_size=18, color=ACCENT_GREEN, bold=True)

    conn_steps = [
        ("User logs in", "Client gets JWT\naccess token", ACCENT_BLUE),
        ("initSocket(token)", "Socket.io connects\nwith auth token", ACCENT_GREEN),
        ("Server Auth", "JWT verified →\nUser loaded from DB", ACCENT_PURPLE),
        ("Auto-Join Rooms", "user:{id} room +\nall conv:{id} rooms", ACCENT_AMBER),
        ("Auto-Deliver", "Mark undelivered msgs\nas delivered → notify", ACCENT_PINK),
        ("Handlers Active", "chatHandler +\npresenceHandler +\ngroupHandler", ACCENT_BLUE),
    ]

    for i, (title, desc, color) in enumerate(conn_steps):
        x = Inches(0.8 + i * 2.05)
        card = add_shape(slide, x, Inches(1.7), Inches(1.8), Inches(1.4), BG_CARD, color)
        add_text_box(slide, x + Inches(0.1), Inches(1.8), Inches(1.6), Inches(0.3),
                     title, font_size=11, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.1), Inches(2.15), Inches(1.6), Inches(0.8),
                     desc, font_size=9, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)
        if i < len(conn_steps) - 1:
            add_right_arrow(slide, x + Inches(1.8), Inches(2.2), Inches(0.25), Inches(0.2), color)

    # Client → Server Events
    add_text_box(slide, Inches(0.8), Inches(3.4), Inches(5.5), Inches(0.4),
                 "📤 Client → Server Events", font_size=16, color=ACCENT_BLUE, bold=True)

    c2s_card = add_shape(slide, Inches(0.8), Inches(3.85), Inches(5.5), Inches(3.3), BG_CARD, ACCENT_BLUE)
    c2s_events = [
        ("send_message", "{ conv_id, content, type, reply_to, expires_in, mentions }"),
        ("typing_start", "{ conversation_id }"),
        ("typing_stop", "{ conversation_id }"),
        ("message_read", "{ message_id?, conversation_id, conversation_type }"),
        ("mark_delivered", "{ message_id }"),
        ("react_message", "{ message_id, emoji }"),
        ("pin_message", "{ message_id, conversation_id }"),
        ("delete_message", "{ message_id, for_all }"),
        ("edit_message", "{ message_id, text }"),
        ("join_conversation", "{ conversation_id, conversation_type }"),
        ("system_message", "{ conversation_id, text }"),
    ]
    for i, (event, payload) in enumerate(c2s_events):
        add_text_box(slide, Inches(1.0), Inches(4.0 + i * 0.28), Inches(1.8), Inches(0.28),
                     event, font_size=9, color=ACCENT_BLUE, bold=True, font_name="Consolas")
        add_text_box(slide, Inches(2.9), Inches(4.0 + i * 0.28), Inches(3.2), Inches(0.28),
                     payload, font_size=8, color=TEXT_GRAY, font_name="Consolas")

    # Server → Client Events
    add_text_box(slide, Inches(6.8), Inches(3.4), Inches(5.5), Inches(0.4),
                 "📥 Server → Client Events", font_size=16, color=ACCENT_GREEN, bold=True)

    s2c_card = add_shape(slide, Inches(6.8), Inches(3.85), Inches(5.5), Inches(3.3), BG_CARD, ACCENT_GREEN)
    s2c_events = [
        ("new_message", "Full message object with sender info"),
        ("message_delivered", "{ message_id, conversation_id }"),
        ("message_read", "{ message_id, reader_id, reader_avatar, read_at }"),
        ("messages_read", "{ message_ids[], reader_id, read_at } (bulk)"),
        ("message_edited", "{ message_id, new_content, edited_at }"),
        ("message_deleted", "{ message_id, deleted_for_all }"),
        ("reaction_update", "{ message_id, reactions[] }"),
        ("user_typing", "{ user_id, display_name, conversation_id }"),
        ("presence_update", "{ user_id, status, last_seen }"),
        ("new_conversation", "Conversation or group object"),
        ("group_updated", "{ group_id, changes }"),
    ]
    for i, (event, payload) in enumerate(s2c_events):
        add_text_box(slide, Inches(7.0), Inches(4.0 + i * 0.28), Inches(1.8), Inches(0.28),
                     event, font_size=9, color=ACCENT_GREEN, bold=True, font_name="Consolas")
        add_text_box(slide, Inches(8.9), Inches(4.0 + i * 0.28), Inches(3.2), Inches(0.28),
                     payload, font_size=8, color=TEXT_GRAY, font_name="Consolas")


def create_chat_workflow_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ACCENT_BLUE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                 "8. One-to-One Chat — Complete Workflow", font_size=32, color=TEXT_WHITE, bold=True)

    # Message Send Flow
    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(12), Inches(0.4),
                 "📨 Sending a Message (end-to-end flow)", font_size=18, color=ACCENT_BLUE, bold=True)

    steps = [
        ("User A types\nmessage", "React MessageInput\ncomponent", ACCENT_BLUE),
        ("socket.emit\n('send_message')", "WebSocket event\nwith content", ACCENT_GREEN),
        ("Server validates\naccess", "resolveConversation\nAccess()", ACCENT_PURPLE),
        ("Friend request\ngate check", "Pending request?\nBlock reply", ACCENT_AMBER),
        ("Message saved\nto MongoDB", "Message.create()\nwith full schema", ACCENT_GREEN),
        ("Check receiver\nonline status", "fetchSockets()\nin user room", ACCENT_PINK),
        ("Broadcast to\nconv room", "socket.to(conv:id)\n.emit('new_message')", ACCENT_BLUE),
        ("Callback to\nsender", "{ success, message }\nwith status", ACCENT_GREEN),
    ]

    for i, (title, desc, color) in enumerate(steps):
        x = Inches(0.4 + i * 1.58)
        card = add_shape(slide, x, Inches(1.7), Inches(1.4), Inches(1.5), BG_CARD, color)
        add_text_box(slide, x + Inches(0.08), Inches(1.8), Inches(1.24), Inches(0.5),
                     title, font_size=10, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.08), Inches(2.3), Inches(1.24), Inches(0.7),
                     desc, font_size=8, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            add_right_arrow(slide, x + Inches(1.4), Inches(2.2), Inches(0.18), Inches(0.15), color)

    # Delivery & Read Receipt Flow
    add_text_box(slide, Inches(0.8), Inches(3.5), Inches(12), Inches(0.4),
                 "✓✓ Delivery & Read Receipt Flow", font_size=18, color=ACCENT_GREEN, bold=True)

    receipt_card = add_shape(slide, Inches(0.8), Inches(4.0), Inches(11.5), Inches(1.8), BG_CARD, BORDER_COLOR)

    receipt_items = [
        ("SENT ✓", "Message saved to MongoDB. Sender sees single checkmark.", ACCENT_BLUE),
        ("DELIVERED ✓✓", "Receiver is online OR comes online later. delivered_to[] updated. Sender sees double gray checkmark.", ACCENT_PURPLE),
        ("READ  ✓✓", "Receiver opens conversation → message_read event → read_receipts[] updated. Sender sees blue double checkmark with reader avatar.", ACCENT_GREEN),
    ]

    for i, (status, desc, color) in enumerate(receipt_items):
        y = Inches(4.15 + i * 0.55)
        add_text_box(slide, Inches(1.0), y, Inches(1.5), Inches(0.4),
                     status, font_size=13, color=color, bold=True)
        add_text_box(slide, Inches(2.6), y, Inches(9.2), Inches(0.5),
                     desc, font_size=11, color=TEXT_GRAY)

    # Features
    add_text_box(slide, Inches(0.8), Inches(6.0), Inches(12), Inches(0.4),
                 "💬 One-to-One Chat Features", font_size=16, color=ACCENT_AMBER, bold=True)

    features = [
        "Text, Image, Video, Audio, File, Sticker, Location",
        "Reply (quote), Forward, Pin messages",
        "Emoji Reactions (👍❤️😂😮😢😡)",
        "Edit messages (within 24h, with history)",
        "Delete: 'just for me' or 'for everyone'",
        "Disappearing messages (1h / 24h / 7d)",
        "Typing indicator ('typing...')",
        "User Block / Unblock / Report",
        "Friend Request gate (accept before chat)",
        "Nicknames for contacts",
    ]

    for i, feat in enumerate(features):
        col = i // 5
        row = i % 5
        add_text_box(slide, Inches(0.8 + col * 6), Inches(6.4 + row * 0.22), Inches(5.5), Inches(0.22),
                     f"• {feat}", font_size=10, color=TEXT_LIGHT_GRAY)


def create_group_chat_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ACCENT_PINK)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                 "9. Group Chat — Feature Workflow", font_size=32, color=TEXT_WHITE, bold=True)

    # Group Creation Flow
    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(12), Inches(0.4),
                 "🏗️ Group Creation & Management", font_size=18, color=ACCENT_PINK, bold=True)

    create_steps = [
        ("Create Group", "Name, description,\navatar, type\n(public/private/secret)", ACCENT_BLUE),
        ("Set Options", "Max members (256),\ninvite_only flag,\nunique_group_id", ACCENT_PURPLE),
        ("Invite Link", "Auto-generated unique\ninvite link + QR code\nfor sharing", ACCENT_GREEN),
        ("Add Members", "Admin adds by user_id\nor members join via\ninvite link", ACCENT_AMBER),
        ("Assign Roles", "Owner → Admin →\nModerator → Member\n(cascading permissions)", ACCENT_PINK),
    ]

    for i, (title, desc, color) in enumerate(create_steps):
        x = Inches(0.8 + i * 2.45)
        card = add_shape(slide, x, Inches(1.7), Inches(2.2), Inches(1.7), BG_CARD, color)
        add_text_box(slide, x + Inches(0.1), Inches(1.8), Inches(2.0), Inches(0.3),
                     title, font_size=13, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.1), Inches(2.15), Inches(2.0), Inches(1.1),
                     desc, font_size=10, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)
        if i < len(create_steps) - 1:
            add_right_arrow(slide, x + Inches(2.2), Inches(2.3), Inches(0.25), Inches(0.2), color)

    # Role Permissions
    add_text_box(slide, Inches(0.8), Inches(3.7), Inches(12), Inches(0.4),
                 "👥 Role-Based Permission System", font_size=18, color=ACCENT_PURPLE, bold=True)

    roles = [
        ("Owner", "🔴", [
            "Delete group", "Change all settings", "Appoint admins",
            "Transfer ownership", "All admin permissions"
        ], ACCENT_RED),
        ("Admin", "🟠", [
            "Add/remove members", "Appoint moderators", "Delete any message",
            "Mute members", "Edit group info"
        ], ACCENT_AMBER),
        ("Moderator", "🟡", [
            "Delete messages", "Mute members", "Review reports",
            "Pin messages", ""
        ], RGBColor(0xFA, 0xCC, 0x15)),
        ("Member", "🟢", [
            "Send messages", "Share files", "React to messages",
            "Leave group", ""
        ], ACCENT_GREEN),
    ]

    for i, (role, emoji, perms, color) in enumerate(roles):
        x = Inches(0.8 + i * 3.1)
        card = add_shape(slide, x, Inches(4.2), Inches(2.8), Inches(2.2), BG_CARD, color)
        add_text_box(slide, x + Inches(0.1), Inches(4.3), Inches(2.6), Inches(0.35),
                     f"{emoji} {role}", font_size=14, color=color, bold=True)
        for j, perm in enumerate(perms):
            if perm:
                add_text_box(slide, x + Inches(0.15), Inches(4.7 + j * 0.28), Inches(2.5), Inches(0.28),
                             f"• {perm}", font_size=10, color=TEXT_GRAY)

    # Group messaging flow
    add_text_box(slide, Inches(0.8), Inches(6.6), Inches(12), Inches(0.4),
                 "💬 Group Message Flow: send_message → check membership → check muted status → Message.create() → broadcast to conv:{group_id} room → all online members receive via Socket.io",
                 font_size=11, color=TEXT_GRAY)


def create_media_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ACCENT_AMBER)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                 "10. File & Media Upload Workflow", font_size=32, color=TEXT_WHITE, bold=True)

    # Upload Flow
    steps = [
        ("User selects file", "React Dropzone\nor file input\n(drag & drop)", ACCENT_BLUE),
        ("Client validation", "Check file type,\nsize limits\n(10-100MB)", ACCENT_PURPLE),
        ("FormData POST", "multipart/form-data\nto /messages/media\nvia Axios", ACCENT_GREEN),
        ("Multer middleware", "Parse multipart\nbody, validate\nMIME type", ACCENT_AMBER),
        ("Cloudinary upload", "cloudinary.uploader\n.upload() → returns\nURL + metadata", ACCENT_PINK),
        ("Save message", "Message.create()\nwith media_url,\nmedia_type, size", ACCENT_BLUE),
        ("Socket broadcast", "Emit 'new_message'\nwith media content\nto conv room", ACCENT_GREEN),
    ]

    for i, (title, desc, color) in enumerate(steps):
        x = Inches(0.4 + i * 1.8)
        card = add_shape(slide, x, Inches(1.3), Inches(1.6), Inches(1.8), BG_CARD, color)
        circle = add_circle(slide, x + Inches(0.55), Inches(1.15), Inches(0.35), color)
        add_text_box(slide, x + Inches(0.6), Inches(1.18), Inches(0.3), Inches(0.3),
                     str(i + 1), font_size=11, color=BG_DARK, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.05), Inches(1.55), Inches(1.5), Inches(0.35),
                     title, font_size=11, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.05), Inches(1.9), Inches(1.5), Inches(0.9),
                     desc, font_size=9, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            add_right_arrow(slide, x + Inches(1.6), Inches(2.0), Inches(0.2), Inches(0.15), color)

    # File Type Limits
    add_text_box(slide, Inches(0.8), Inches(3.5), Inches(12), Inches(0.4),
                 "📁 Supported File Types & Limits", font_size=18, color=ACCENT_AMBER, bold=True)

    file_types = [
        ("📷 Images", "10 MB", "JPG, PNG, GIF, WEBP", "Cloudinary auto-compress", ACCENT_BLUE),
        ("🎬 Videos", "100 MB", "MP4, MOV, AVI", "Cloudinary HLS streaming", ACCENT_PURPLE),
        ("🎵 Audio", "25 MB", "MP3, OGG, WAV, M4A", "Cloudinary storage", ACCENT_GREEN),
        ("📄 Documents", "50 MB", "PDF, DOCX, XLSX, TXT", "Cloudinary storage", ACCENT_AMBER),
        ("🎙️ Voice Messages", "5 min", "OGG / WEBM", "Cloudinary storage", ACCENT_PINK),
    ]

    for i, (type_name, limit, formats, storage, color) in enumerate(file_types):
        x = Inches(0.8 + i * 2.45)
        card = add_shape(slide, x, Inches(4.0), Inches(2.2), Inches(1.8), BG_CARD, color)
        add_text_box(slide, x + Inches(0.1), Inches(4.1), Inches(2.0), Inches(0.3),
                     type_name, font_size=13, color=color, bold=True)
        add_text_box(slide, x + Inches(0.1), Inches(4.4), Inches(2.0), Inches(0.3),
                     f"Max: {limit}", font_size=11, color=TEXT_WHITE)
        add_text_box(slide, x + Inches(0.1), Inches(4.7), Inches(2.0), Inches(0.3),
                     formats, font_size=10, color=TEXT_GRAY)
        add_text_box(slide, x + Inches(0.1), Inches(5.0), Inches(2.0), Inches(0.5),
                     storage, font_size=9, color=TEXT_GRAY)

    # Cloudinary
    add_text_box(slide, Inches(0.8), Inches(6.0), Inches(12), Inches(0.4),
                 "☁️ Cloudinary Integration", font_size=18, color=ACCENT_BLUE, bold=True)

    cloud_card = add_shape(slide, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.8), BG_CARD, BORDER_COLOR)
    add_text_box(slide, Inches(1.0), Inches(6.5), Inches(11), Inches(0.6),
                 "Free tier: 25GB storage  •  Multer middleware handles multipart parsing  •  "
                 "multer-storage-cloudinary SDK for direct upload  •  Returns secure_url for CDN access  •  "
                 "Auto-generates thumbnails for images & videos  •  Supports transformations (resize, crop, compress)",
                 font_size=11, color=TEXT_GRAY)


def create_message_lifecycle_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ACCENT_INDIGO)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                 "11. Message Lifecycle — Edit, Delete, Pin, React", font_size=32, color=TEXT_WHITE, bold=True)

    # Edit
    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(5.8), Inches(0.4),
                 "✏️ Edit Message", font_size=18, color=ACCENT_BLUE, bold=True)
    edit_card = add_shape(slide, Inches(0.8), Inches(1.6), Inches(5.8), Inches(1.5), BG_CARD, ACCENT_BLUE)
    edit_items = [
        "1. User emits 'edit_message' { message_id, text }",
        "2. Server checks: sender owns message? Within 24 hours?",
        "3. Old content pushed to edit_history[] array",
        "4. content.text updated, is_edited = true, edited_at set",
        "5. 'message_edited' broadcast to conv room with new content",
        "⚠️ Messages older than 24h cannot be edited",
    ]
    for i, item in enumerate(edit_items):
        add_text_box(slide, Inches(1.0), Inches(1.7 + i * 0.22), Inches(5.4), Inches(0.22),
                     item, font_size=10, color=TEXT_GRAY)

    # Delete
    add_text_box(slide, Inches(7.0), Inches(1.2), Inches(5.8), Inches(0.4),
                 "🗑️ Delete Message", font_size=18, color=ACCENT_RED, bold=True)
    del_card = add_shape(slide, Inches(7.0), Inches(1.6), Inches(5.8), Inches(1.5), BG_CARD, ACCENT_RED)
    del_items = [
        "1. User emits 'delete_message' { message_id, for_all }",
        "2. If for_all AND sender is message owner:",
        "   → deleted_for_all = true, is_deleted = true",
        "3. If NOT for_all (delete for me only):",
        "   → user_id added to deleted_for[] array",
        "4. 'message_deleted' broadcast (for_all) or sent to self only",
    ]
    for i, item in enumerate(del_items):
        add_text_box(slide, Inches(7.2), Inches(1.7 + i * 0.22), Inches(5.4), Inches(0.22),
                     item, font_size=10, color=TEXT_GRAY)

    # Pin
    add_text_box(slide, Inches(0.8), Inches(3.4), Inches(5.8), Inches(0.4),
                 "📌 Pin Message", font_size=18, color=ACCENT_AMBER, bold=True)
    pin_card = add_shape(slide, Inches(0.8), Inches(3.8), Inches(5.8), Inches(1.2), BG_CARD, ACCENT_AMBER)
    pin_items = [
        "1. User emits 'pin_message' { message_id, conversation_id }",
        "2. Server toggles is_pinned (pin/unpin)",
        "3. 'message_pinned' broadcast to conv room",
        "4. UI shows pinned messages at top / in info panel",
    ]
    for i, item in enumerate(pin_items):
        add_text_box(slide, Inches(1.0), Inches(3.9 + i * 0.25), Inches(5.4), Inches(0.25),
                     item, font_size=10, color=TEXT_GRAY)

    # React
    add_text_box(slide, Inches(7.0), Inches(3.4), Inches(5.8), Inches(0.4),
                 "😊 Emoji Reactions", font_size=18, color=ACCENT_PINK, bold=True)
    react_card = add_shape(slide, Inches(7.0), Inches(3.8), Inches(5.8), Inches(1.2), BG_CARD, ACCENT_PINK)
    react_items = [
        "1. User emits 'react_message' { message_id, emoji }",
        "2. Server removes user's old reaction (if any)",
        "3. Adds new { user_id, emoji, reacted_at } to reactions[]",
        "4. 'reaction_update' broadcast with full reactions array",
    ]
    for i, item in enumerate(react_items):
        add_text_box(slide, Inches(7.2), Inches(3.9 + i * 0.25), Inches(5.4), Inches(0.25),
                     item, font_size=10, color=TEXT_GRAY)

    # Disappearing Messages
    add_text_box(slide, Inches(0.8), Inches(5.3), Inches(5.8), Inches(0.4),
                 "⏱️ Disappearing Messages", font_size=18, color=ACCENT_GREEN, bold=True)
    dis_card = add_shape(slide, Inches(0.8), Inches(5.7), Inches(5.8), Inches(1.5), BG_CARD, ACCENT_GREEN)
    dis_items = [
        "1. Sender sets expires_in when sending ('1h', '24h', '7d')",
        "2. Server calculates expires_at = now + duration",
        "3. MongoDB TTL index: messageSchema.index({ expires_at: 1 },",
        "   { expireAfterSeconds: 0 })",
        "4. MongoDB automatically deletes when expires_at is reached",
        "5. No cron job needed — database handles cleanup",
    ]
    for i, item in enumerate(dis_items):
        add_text_box(slide, Inches(1.0), Inches(5.8 + i * 0.22), Inches(5.4), Inches(0.22),
                     item, font_size=10, color=TEXT_GRAY)

    # Reply & Forward
    add_text_box(slide, Inches(7.0), Inches(5.3), Inches(5.8), Inches(0.4),
                 "↩️ Reply & Forward", font_size=18, color=ACCENT_PURPLE, bold=True)
    reply_card = add_shape(slide, Inches(7.0), Inches(5.7), Inches(5.8), Inches(1.5), BG_CARD, ACCENT_PURPLE)
    reply_items = [
        "Reply: reply_to field → ObjectId ref to original Message",
        "  → Populated with original message content + sender info",
        "  → ReplyPreview component shows quoted message above input",
        "Forward: forwarded_from field → ObjectId ref to source",
        "  → Can forward to multiple conversations",
        "  → Original sender info preserved",
    ]
    for i, item in enumerate(reply_items):
        add_text_box(slide, Inches(7.2), Inches(5.8 + i * 0.22), Inches(5.4), Inches(0.22),
                     item, font_size=10, color=TEXT_GRAY)


def create_notification_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ACCENT_PINK)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                 "12. Notification & Presence System", font_size=32, color=TEXT_WHITE, bold=True)

    # Presence System
    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(5.8), Inches(0.4),
                 "🟢 Online Presence System", font_size=18, color=ACCENT_GREEN, bold=True)

    pres_card = add_shape(slide, Inches(0.8), Inches(1.7), Inches(5.8), Inches(2.8), BG_CARD, ACCENT_GREEN)
    pres_items = [
        ("Connect", "Socket connects → joins user:{id} room → 'online' status broadcast to all users", ACCENT_GREEN),
        ("Disconnect", "Socket disconnects → check remaining sockets for same user → if none, update last_seen in DB → broadcast 'offline'", ACCENT_RED),
        ("Multi-tab", "Multiple tabs share user room → offline only when ALL tabs close", ACCENT_BLUE),
        ("Status", "Users can set: online / away / busy / invisible → custom status message → visibility settings (everyone/contacts/nobody)", ACCENT_PURPLE),
    ]

    for i, (title, desc, color) in enumerate(pres_items):
        add_text_box(slide, Inches(1.0), Inches(1.85 + i * 0.65), Inches(1.3), Inches(0.3),
                     title, font_size=12, color=color, bold=True)
        add_text_box(slide, Inches(2.3), Inches(1.85 + i * 0.65), Inches(4.0), Inches(0.6),
                     desc, font_size=10, color=TEXT_GRAY)

    # Typing Indicators
    add_text_box(slide, Inches(7.0), Inches(1.2), Inches(5.8), Inches(0.4),
                 "⌨️ Typing Indicators", font_size=18, color=ACCENT_BLUE, bold=True)

    type_card = add_shape(slide, Inches(7.0), Inches(1.7), Inches(5.8), Inches(1.2), BG_CARD, ACCENT_BLUE)
    type_items = [
        "1. User starts typing → emit 'typing_start' { conversation_id }",
        "2. Server broadcasts 'user_typing' to conv room (except sender)",
        "3. User stops typing → emit 'typing_stop' { conversation_id }",
        "4. Chat store tracks typingUsers per conversation",
    ]
    for i, item in enumerate(type_items):
        add_text_box(slide, Inches(7.2), Inches(1.8 + i * 0.25), Inches(5.4), Inches(0.25),
                     item, font_size=10, color=TEXT_GRAY)

    # Unread Counts
    add_text_box(slide, Inches(7.0), Inches(3.2), Inches(5.8), Inches(0.4),
                 "🔔 Unread Count System", font_size=18, color=ACCENT_AMBER, bold=True)

    unread_card = add_shape(slide, Inches(7.0), Inches(3.6), Inches(5.8), Inches(1.0), BG_CARD, ACCENT_AMBER)
    unread_items = [
        "• Zustand store: unreadCounts object { convId: number }",
        "• Increment on new_message (if not active conversation)",
        "• Clear when user opens/selects the conversation",
    ]
    for i, item in enumerate(unread_items):
        add_text_box(slide, Inches(7.2), Inches(3.7 + i * 0.25), Inches(5.4), Inches(0.25),
                     item, font_size=10, color=TEXT_GRAY)

    # Email Notifications
    add_text_box(slide, Inches(0.8), Inches(4.8), Inches(12), Inches(0.4),
                 "📧 Email Notification System (Brevo API)", font_size=18, color=ACCENT_PINK, bold=True)

    email_card = add_shape(slide, Inches(0.8), Inches(5.3), Inches(11.5), Inches(2.0), BG_CARD, BORDER_COLOR)

    email_flow = [
        ("OTP Generation", "crypto.randomInt() → 6-digit code → stored in Redis (key: otp:{type}:{email}, TTL: 10 min) or in-memory Map fallback", ACCENT_BLUE),
        ("Email Sending", "Brevo REST API (api.brevo.com/v3/smtp/email) → HTML template with gradient header, styled OTP code, expiry info", ACCENT_GREEN),
        ("OTP Verification", "Compare submitted code with stored code → one-time use (deleted after verify) → handles Redis failure with memory fallback", ACCENT_PURPLE),
        ("Use Cases", "Email verification on registration  •  Password reset  •  Resend OTP endpoint available", ACCENT_AMBER),
    ]

    for i, (title, desc, color) in enumerate(email_flow):
        add_text_box(slide, Inches(1.0), Inches(5.4 + i * 0.47), Inches(1.5), Inches(0.3),
                     title, font_size=11, color=color, bold=True)
        add_text_box(slide, Inches(2.5), Inches(5.4 + i * 0.47), Inches(9.5), Inches(0.45),
                     desc, font_size=10, color=TEXT_GRAY)


def create_frontend_arch_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ACCENT_BLUE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                 "13. Frontend Architecture & Component Tree", font_size=32, color=TEXT_WHITE, bold=True)

    # Pages/Routes
    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(3.5), Inches(0.4),
                 "🗺️ Routes (App.jsx)", font_size=16, color=ACCENT_GREEN, bold=True)
    routes_card = add_shape(slide, Inches(0.8), Inches(1.6), Inches(3.5), Inches(2.0), BG_CARD, ACCENT_GREEN)
    routes = [
        "/              → LandingPage",
        "/login          → LoginPage (Public)",
        "/register       → RegisterPage (Public)",
        "/forgot-password → ForgotPasswordPage",
        "/app            → ChatLayout (Protected)",
        "/*              → Redirect to /",
    ]
    for i, route in enumerate(routes):
        add_text_box(slide, Inches(1.0), Inches(1.7 + i * 0.28), Inches(3.1), Inches(0.28),
                     route, font_size=10, color=TEXT_GRAY, font_name="Consolas")

    # Component Tree
    add_text_box(slide, Inches(4.8), Inches(1.2), Inches(8), Inches(0.4),
                 "🌳 Component Tree (ChatLayout)", font_size=16, color=ACCENT_PURPLE, bold=True)

    tree_card = add_shape(slide, Inches(4.8), Inches(1.6), Inches(7.8), Inches(5.5), BG_CARD, ACCENT_PURPLE)

    # Three-column layout description
    columns = [
        ("Sidebar (280px)", ACCENT_BLUE, [
            "Sidebar.jsx",
            "├── User profile header",
            "├── Tab bar (Chats/Groups/Contacts)",
            "├── Search input",
            "├── ConversationList.jsx",
            "│   └── Conversation items",
            "├── GroupList.jsx",
            "│   └── Group items",
            "└── ContactSearch.jsx",
            "    └── User search results",
        ]),
        ("ChatWindow (flex)", ACCENT_GREEN, [
            "ChatWindow.jsx",
            "├── ChatHeader.jsx",
            "│   └── User/group info, actions",
            "├── Messages area (virtualized)",
            "│   ├── DateSeparator.jsx",
            "│   ├── MessageItem.jsx",
            "│   │   ├── Message bubble",
            "│   │   ├── ReactionPicker.jsx",
            "│   │   └── Reply/Forward/Edit UIs",
            "│   └── EmptyChatState.jsx",
            "├── ReplyPreview.jsx",
            "└── MessageInput.jsx",
            "    └── Emoji, attachments, send",
        ]),
        ("InfoPanel (320px)", ACCENT_AMBER, [
            "InfoPanel.jsx (43KB)",
            "├── User/Group profile view",
            "├── Shared media gallery",
            "├── Member list (groups)",
            "├── Role management",
            "├── Nickname settings",
            "├── Block/Report actions",
            "├── Group settings",
            "└── Leave/Delete options",
            "",
            "Modals:",
            "├── NewChatModal.jsx",
            "└── UserSettingsModal.jsx",
        ]),
    ]

    for i, (title, color, items) in enumerate(columns):
        x = Inches(5.0 + i * 2.5)
        add_text_box(slide, x, Inches(1.7), Inches(2.3), Inches(0.3),
                     title, font_size=11, color=color, bold=True)
        for j, item in enumerate(items):
            add_text_box(slide, x, Inches(2.05 + j * 0.35), Inches(2.3), Inches(0.35),
                         item, font_size=8, color=TEXT_GRAY, font_name="Consolas")

    # State Management
    add_text_box(slide, Inches(0.8), Inches(3.9), Inches(3.5), Inches(0.4),
                 "🗄️ State Management", font_size=16, color=ACCENT_AMBER, bold=True)

    store_card = add_shape(slide, Inches(0.8), Inches(4.3), Inches(3.5), Inches(2.8), BG_CARD, ACCENT_AMBER)
    stores = [
        ("authStore.js (Zustand + persist)", ACCENT_BLUE, [
            "user, tokens, isAuthenticated",
            "login(), register(), logout()",
            "fetchMe(), updateUser()",
            "Persisted to localStorage",
        ]),
        ("chatStore.js (Zustand)", ACCENT_GREEN, [
            "activeConversation, conversations[]",
            "groups[], messages{}, unreadCounts{}",
            "typingUsers{}, presenceMap{}",
            "replyTo, searchQuery, sidebarTab",
            "Pinned conversations/groups",
            "Nickname management",
        ]),
    ]

    y_offset = Inches(4.4)
    for store_name, color, items in stores:
        add_text_box(slide, Inches(1.0), y_offset, Inches(3.1), Inches(0.3),
                     store_name, font_size=10, color=color, bold=True)
        y_offset += Inches(0.3)
        for item in items:
            add_text_box(slide, Inches(1.1), y_offset, Inches(3.0), Inches(0.2),
                         f"• {item}", font_size=8, color=TEXT_GRAY)
            y_offset += Inches(0.2)
        y_offset += Inches(0.15)


def create_api_summary_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ACCENT_GREEN)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                 "14. API Endpoint Summary", font_size=32, color=TEXT_WHITE, bold=True)

    api_groups = [
        ("🔐 Auth API", ACCENT_BLUE, [
            "POST /api/auth/register — Create account",
            "POST /api/auth/login — Login + JWT",
            "POST /api/auth/logout — Logout + blacklist",
            "POST /api/auth/refresh — Rotate tokens",
            "POST /api/auth/forgot-password — OTP email",
            "POST /api/auth/reset-password — New password",
            "POST /api/auth/verify-email — Confirm OTP",
            "POST /api/auth/2fa/enable — Setup 2FA",
            "POST /api/auth/2fa/verify — Confirm 2FA",
            "POST /api/auth/resend-otp — Resend code",
        ]),
        ("👤 User API", ACCENT_GREEN, [
            "GET    /api/users/me — My profile",
            "PUT    /api/users/me — Update profile",
            "POST   /api/users/me/avatar — Upload avatar",
            "GET    /api/users/search?q= — Search users",
            "GET    /api/users/:id — View user",
            "POST   /api/users/:id/block — Block user",
            "DELETE /api/users/:id/block — Unblock",
            "POST   /api/users/:id/report — Report user",
        ]),
        ("💬 Conversation API", ACCENT_PURPLE, [
            "GET    /api/conversations — List all",
            "POST   /api/conversations — Start new chat",
            "GET    /api/conversations/:id — Get details",
            "DELETE /api/conversations/:id — Delete chat",
            "GET    /api/conversations/:id/messages — Messages",
            "POST   /api/conversations/:id/messages — Send msg",
            "PUT    /api/conversations/:id/nickname — Set nick",
            "POST   /api/conversations/:id/accept-request",
        ]),
        ("👥 Group API", ACCENT_AMBER, [
            "GET    /api/groups — My groups",
            "POST   /api/groups — Create group",
            "GET    /api/groups/:id — Group details",
            "PUT    /api/groups/:id — Update group",
            "DELETE /api/groups/:id — Delete group",
            "GET    /api/groups/:id/members — Members",
            "POST   /api/groups/:id/members — Add member",
            "PUT    /api/groups/:id/members/:uid/role",
            "POST   /api/groups/join/:invite — Join via link",
            "POST   /api/groups/:id/leave — Leave group",
        ]),
    ]

    for i, (title, color, endpoints) in enumerate(api_groups):
        col = i % 2
        row = i // 2
        x = Inches(0.8 + col * 6.2)
        y = Inches(1.2 + row * 3.3)

        add_text_box(slide, x, y, Inches(5.8), Inches(0.4),
                     title, font_size=15, color=color, bold=True)

        card = add_shape(slide, x, y + Inches(0.35), Inches(5.8), Inches(2.7), BG_CARD, color)

        for j, ep in enumerate(endpoints):
            add_text_box(slide, x + Inches(0.15), y + Inches(0.45 + j * 0.25), Inches(5.5), Inches(0.25),
                         ep, font_size=9, color=TEXT_GRAY, font_name="Consolas")


def create_deployment_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ACCENT_GREEN)

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
                 "15. Deployment — Server Architecture", font_size=30, color=TEXT_WHITE, bold=True)

    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.35),
                 "Investigated LIVE from DigitalOcean Droplet  |  tongchat.app  |  Ubuntu 22.04  |  $6/month",
                 font_size=13, color=TEXT_GRAY)

    # Important note
    note_card = add_shape(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.55), RGBColor(0x1C, 0x1C, 0x0A), ACCENT_AMBER)
    add_text_box(slide, Inches(1.0), Inches(1.35), Inches(11.0), Inches(0.45),
                 "NOT DOCKER: Despite docker-compose.yml in repo, the droplet runs natively with Nginx + PM2 + Node.js (no Docker installed)",
                 font_size=12, color=ACCENT_AMBER, bold=True)

    # === Request Flow ===
    add_text_box(slide, Inches(0.8), Inches(2.0), Inches(12), Inches(0.35),
                 "Request Flow: User -> Nginx -> Backend", font_size=16, color=ACCENT_BLUE, bold=True)

    # Flow: User -> Nginx -> Route Decision
    flow_items = [
        ("User Browser", "https://tongchat.app\n(HTTPS request)", ACCENT_BLUE),
        ("Nginx :443", "SSL termination\n(Let's Encrypt cert)", ACCENT_GREEN),
        ("Route Decision", "Static file?\nAPI call? WebSocket?", ACCENT_PURPLE),
    ]
    for i, (title, desc, color) in enumerate(flow_items):
        x = Inches(0.8 + i * 2.3)
        card = add_shape(slide, x, Inches(2.4), Inches(2.0), Inches(1.1), BG_CARD, color)
        add_text_box(slide, x + Inches(0.1), Inches(2.45), Inches(1.8), Inches(0.3),
                     title, font_size=11, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.1), Inches(2.75), Inches(1.8), Inches(0.6),
                     desc, font_size=9, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)
        if i < 2:
            add_right_arrow(slide, x + Inches(2.0), Inches(2.8), Inches(0.3), Inches(0.2), color)

    # Three branches
    branches = [
        ("Static Files", "/ , /login, /app\n/var/www/tong/\nfrontend/dist/\nindex.html + assets", ACCENT_BLUE, Inches(5.8)),
        ("/api/* Proxy", "proxy_pass to\n127.0.0.1:5000\nExpress.js REST API", ACCENT_GREEN, Inches(8.0)),
        ("/socket.io/* WS", "WebSocket upgrade\nproxy to :5000\nConnection: upgrade", ACCENT_PURPLE, Inches(10.2)),
    ]
    for title, desc, color, x in branches:
        add_arrow(slide, x + Inches(0.6), Inches(3.5), Inches(0.25), Inches(0.3), color)
        card = add_shape(slide, x, Inches(3.9), Inches(1.8), Inches(1.2), BG_CARD, color)
        add_text_box(slide, x + Inches(0.05), Inches(3.95), Inches(1.7), Inches(0.25),
                     title, font_size=10, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.05), Inches(4.2), Inches(1.7), Inches(0.8),
                     desc, font_size=8, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER, font_name="Consolas")

    # PM2 + Node.js Backend
    add_arrow(slide, Inches(9.3), Inches(5.1), Inches(0.25), Inches(0.25), ACCENT_GREEN)
    pm2_card = add_shape(slide, Inches(7.5), Inches(5.4), Inches(4.8), Inches(1.2), BG_CARD, ACCENT_GREEN)
    add_text_box(slide, Inches(7.65), Inches(5.45), Inches(4.5), Inches(0.3),
                 "PM2 v6.0.14 -> Node.js v20.20.2 (port 5000)", font_size=12, color=ACCENT_GREEN, bold=True)
    pm2_items = [
        "Process: tong-backend  |  Status: Online  |  Uptime: 32h+",
        "Script: /var/www/tong/backend/src/server.js",
        "Memory: ~97 MB  |  Restarts: 0  |  Mode: fork",
    ]
    for i, item in enumerate(pm2_items):
        add_text_box(slide, Inches(7.65), Inches(5.8 + i * 0.22), Inches(4.5), Inches(0.22),
                     item, font_size=9, color=TEXT_GRAY, font_name="Consolas")

    # Droplet Specs
    add_text_box(slide, Inches(0.8), Inches(5.2), Inches(6), Inches(0.35),
                 "Droplet Specs & Nginx Config", font_size=14, color=ACCENT_AMBER, bold=True)

    specs_card = add_shape(slide, Inches(0.8), Inches(5.55), Inches(3.0), Inches(1.7), BG_CARD, ACCENT_AMBER)
    specs = [
        ("IP Address", "168.144.78.31"),
        ("OS", "Ubuntu 22.04 LTS"),
        ("Plan", "$6/month droplet"),
        ("Disk", "25 GB (3.4 GB used)"),
        ("RAM", "957 MB (265 MB used)"),
        ("Firewall", "UFW inactive"),
    ]
    for i, (label, value) in enumerate(specs):
        add_text_box(slide, Inches(0.95), Inches(5.6 + i * 0.26), Inches(1.2), Inches(0.26),
                     label, font_size=9, color=ACCENT_AMBER, bold=True)
        add_text_box(slide, Inches(2.1), Inches(5.6 + i * 0.26), Inches(1.5), Inches(0.26),
                     value, font_size=9, color=TEXT_GRAY)

    # SSL info
    ssl_card = add_shape(slide, Inches(4.0), Inches(5.55), Inches(3.3), Inches(1.7), BG_CARD, ACCENT_BLUE)
    add_text_box(slide, Inches(4.15), Inches(5.6), Inches(3.0), Inches(0.25),
                 "SSL / Domain", font_size=11, color=ACCENT_BLUE, bold=True)
    ssl_items = [
        "Provider: Let's Encrypt (Certbot)",
        "Domain: tongchat.app + www",
        "Expiry: July 8, 2026 (auto-renew)",
        "HTTP -> HTTPS redirect (port 80)",
        "Ports: 22 (SSH), 80, 443, 5000",
    ]
    for i, item in enumerate(ssl_items):
        add_text_box(slide, Inches(4.15), Inches(5.9 + i * 0.24), Inches(3.0), Inches(0.24),
                     item, font_size=9, color=TEXT_GRAY)


def create_deployment_flow_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ACCENT_AMBER)

    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
                 "16. Deployment — Deploy Flow & External Services", font_size=30, color=TEXT_WHITE, bold=True)

    # === Deploy Script ===
    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(6), Inches(0.35),
                 "deploy.sh — How Updates Are Deployed", font_size=16, color=ACCENT_GREEN, bold=True)

    # Deploy flow steps
    deploy_steps = [
        ("1", "Developer\npushes code", "git push\norigin main", ACCENT_BLUE),
        ("2", "SSH into\ndroplet", "ssh root@\n168.144.78.31", ACCENT_GREEN),
        ("3", "Pull latest\ncode", "git pull\norigin main", ACCENT_PURPLE),
        ("4", "Update\nbackend", "npm install\npm2 restart", ACCENT_AMBER),
        ("5", "Build\nfrontend", "npm install\nnpm run build", ACCENT_PINK),
        ("6", "Live!\nNginx serves", "New dist/ files\nauto-served", ACCENT_GREEN),
    ]

    for i, (num, title, desc, color) in enumerate(deploy_steps):
        x = Inches(0.8 + i * 2.05)
        card = add_shape(slide, x, Inches(1.55), Inches(1.8), Inches(1.5), BG_CARD, color)
        circle = add_circle(slide, x + Inches(0.65), Inches(1.4), Inches(0.35), color)
        add_text_box(slide, x + Inches(0.73), Inches(1.43), Inches(0.2), Inches(0.3),
                     num, font_size=12, color=BG_DARK, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.05), Inches(1.85), Inches(1.7), Inches(0.45),
                     title, font_size=10, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.05), Inches(2.3), Inches(1.7), Inches(0.6),
                     desc, font_size=9, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER, font_name="Consolas")
        if i < len(deploy_steps) - 1:
            add_right_arrow(slide, x + Inches(1.8), Inches(2.1), Inches(0.25), Inches(0.2), color)

    # Deploy script content
    add_text_box(slide, Inches(0.8), Inches(3.3), Inches(6), Inches(0.35),
                 "deploy.sh (actual script on server)", font_size=13, color=ACCENT_AMBER, bold=True)
    script_card = add_shape(slide, Inches(0.8), Inches(3.65), Inches(5.8), Inches(1.8), BG_CARD, ACCENT_AMBER)
    script_lines = [
        '#!/bin/bash',
        'git pull origin main',
        '',
        'cd backend',
        'npm install --legacy-peer-deps',
        'pm2 restart tong-backend',
        'cd ..',
        '',
        'cd frontend',
        'npm install --legacy-peer-deps',
        'npm run build          # Vite -> dist/',
        'cd ..',
    ]
    for i, line in enumerate(script_lines):
        add_text_box(slide, Inches(1.0), Inches(3.7 + i * 0.145), Inches(5.4), Inches(0.145),
                     line, font_size=8, color=TEXT_GRAY if line else TEXT_GRAY, font_name="Consolas")

    # Useful commands
    add_text_box(slide, Inches(0.8), Inches(5.6), Inches(6), Inches(0.3),
                 "Useful Server Commands", font_size=13, color=ACCENT_BLUE, bold=True)
    cmds_card = add_shape(slide, Inches(0.8), Inches(5.9), Inches(5.8), Inches(1.3), BG_CARD, ACCENT_BLUE)
    cmds = [
        "pm2 list                   # See all processes",
        "pm2 logs tong-backend      # Stream live logs",
        "pm2 restart tong-backend   # Restart backend",
        "pm2 monit                  # Real-time CPU/memory",
        "nginx -t && systemctl reload nginx  # Test & reload",
        "certbot renew              # Renew SSL certificate",
    ]
    for i, cmd in enumerate(cmds):
        add_text_box(slide, Inches(1.0), Inches(5.95 + i * 0.2), Inches(5.4), Inches(0.2),
                     cmd, font_size=8, color=TEXT_GRAY, font_name="Consolas")

    # === External Services ===
    add_text_box(slide, Inches(7.0), Inches(3.3), Inches(6), Inches(0.35),
                 "External Cloud Services (Free Tier)", font_size=14, color=ACCENT_PURPLE, bold=True)

    ext_services = [
        ("PostgreSQL", "Supabase", "Users, Groups,\nConversations, Auth", "500 MB", ACCENT_PURPLE),
        ("MongoDB", "Atlas M0", "Messages, Reactions,\nRead Receipts", "512 MB", ACCENT_GREEN),
        ("Redis", "Redis Cloud", "OTP codes, Token\nblacklist, cache", "10K req/day", ACCENT_RED),
        ("File Storage", "Cloudinary", "Images, Videos,\nAudio, Documents", "25 GB", ACCENT_AMBER),
        ("Email", "Brevo API", "OTP verification\nemails (HTML)", "300/day", ACCENT_PINK),
    ]

    for i, (name, provider, desc, limit, color) in enumerate(ext_services):
        y = Inches(3.7 + i * 0.82)
        card = add_shape(slide, Inches(7.0), y, Inches(5.5), Inches(0.72), BG_CARD, color)
        add_text_box(slide, Inches(7.15), y + Inches(0.05), Inches(1.5), Inches(0.25),
                     name, font_size=11, color=color, bold=True)
        add_text_box(slide, Inches(8.6), y + Inches(0.05), Inches(1.2), Inches(0.25),
                     provider, font_size=10, color=TEXT_WHITE, bold=True)
        add_text_box(slide, Inches(8.6), y + Inches(0.32), Inches(2.0), Inches(0.35),
                     desc, font_size=8, color=TEXT_GRAY)
        add_text_box(slide, Inches(10.8), y + Inches(0.15), Inches(1.5), Inches(0.3),
                     f"Free: {limit}", font_size=9, color=ACCENT_GREEN, bold=True)

    # Note
    add_text_box(slide, Inches(7.0), Inches(7.0), Inches(5.5), Inches(0.3),
                 "All databases are external — NOT hosted on the droplet",
                 font_size=10, color=TEXT_GRAY)


def create_file_structure_slide(prs):  # Slide 17
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ACCENT_INDIGO)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                 "17. Project File Structure", font_size=32, color=TEXT_WHITE, bold=True)

    # Backend
    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(5.8), Inches(0.4),
                 "⚙️ Backend (backend/src/)", font_size=16, color=ACCENT_GREEN, bold=True)

    be_card = add_shape(slide, Inches(0.8), Inches(1.6), Inches(5.8), Inches(5.5), BG_CARD, ACCENT_GREEN)
    be_files = [
        "server.js                  — Entry point (HTTP + Socket.io)",
        "app.js                     — Express app (middleware + routes)",
        "",
        "config/",
        "  ├── env.js               — Environment variable loader",
        "  ├── database.js           — PostgreSQL + MongoDB + Redis",
        "  └── cloudinary.js         — Cloudinary SDK config",
        "",
        "controllers/",
        "  ├── auth.controller.js    — Register, Login, 2FA, OTP",
        "  ├── user.controller.js    — Profile, Search, Block, Report",
        "  ├── conversation.controller.js — Direct chat CRUD",
        "  ├── message.controller.js — Message edit, delete, search",
        "  └── group.controller.js   — Group CRUD, members, roles",
        "",
        "routes/   (5 route files mirror controllers)",
        "",
        "middleware/",
        "  ├── auth.js               — JWT verification middleware",
        "  ├── rateLimiter.js         — Rate limiting (100 req/min)",
        "  ├── upload.js              — Multer + Cloudinary upload",
        "  ├── validate.js            — express-validator wrapper",
        "  └── errorHandler.js        — Global error handler",
        "",
        "services/",
        "  ├── token.service.js       — JWT generate/verify/blacklist",
        "  ├── otp.service.js         — OTP generate/store/verify",
        "  ├── email.service.js       — Brevo email sending",
        "  └── cloudinary.service.js  — File upload helper",
        "",
        "socket/",
        "  ├── index.js               — Socket.io init + auth + rooms",
        "  ├── chatHandler.js         — Message events (send/edit/del)",
        "  ├── presenceHandler.js     — Online status events",
        "  └── groupHandler.js        — Group-specific events",
        "",
        "models/",
        "  ├── Message.js             — MongoDB message schema",
        "  └── ConversationVisibility.js",
        "",
        "utils/   (ApiError, ApiResponse, helpers)",
    ]

    for i, line in enumerate(be_files):
        add_text_box(slide, Inches(1.0), Inches(1.65 + i * 0.148), Inches(5.4), Inches(0.15),
                     line, font_size=7, color=TEXT_GRAY if line else TEXT_GRAY, font_name="Consolas")

    # Frontend
    add_text_box(slide, Inches(7.0), Inches(1.2), Inches(5.8), Inches(0.4),
                 "🖥️ Frontend (frontend/src/)", font_size=16, color=ACCENT_BLUE, bold=True)

    fe_card = add_shape(slide, Inches(7.0), Inches(1.6), Inches(5.8), Inches(5.5), BG_CARD, ACCENT_BLUE)
    fe_files = [
        "main.jsx                   — React app entry point",
        "App.jsx                    — Router + Protected/Public routes",
        "index.css                  — Global styles (TailwindCSS)",
        "",
        "pages/",
        "  ├── LandingPage.jsx       — Public landing page",
        "  ├── ChatLayout.jsx        — Main chat 3-column layout",
        "  └── auth/",
        "      ├── LoginPage.jsx     — Login form + 2FA",
        "      ├── RegisterPage.jsx  — Registration + OTP verify",
        "      └── ForgotPasswordPage.jsx",
        "",
        "components/",
        "  ├── sidebar/",
        "  │   ├── Sidebar.jsx       — Tab bar + search",
        "  │   ├── ConversationList.jsx — Chat list",
        "  │   ├── GroupList.jsx      — Group list",
        "  │   └── ContactSearch.jsx  — User search",
        "  ├── chat/",
        "  │   ├── ChatWindow.jsx     — Message viewport",
        "  │   ├── ChatHeader.jsx     — Chat top bar",
        "  │   ├── MessageItem.jsx    — Message bubble (24KB)",
        "  │   ├── MessageInput.jsx   — Input bar (18KB)",
        "  │   ├── ReplyPreview.jsx   — Reply quote",
        "  │   ├── ReactionPicker.jsx — Emoji picker",
        "  │   ├── DateSeparator.jsx",
        "  │   └── EmptyChatState.jsx",
        "  ├── infopanel/",
        "  │   └── InfoPanel.jsx      — Details panel (43KB)",
        "  └── modals/",
        "      ├── NewChatModal.jsx   — New chat/group",
        "      └── UserSettingsModal.jsx",
        "",
        "store/   (authStore, chatStore)",
        "hooks/   (useSeo, useSocketEvents)",
        "lib/     (api, apiServices, socket)",
    ]

    for i, line in enumerate(fe_files):
        add_text_box(slide, Inches(7.2), Inches(1.65 + i * 0.148), Inches(5.4), Inches(0.15),
                     line, font_size=7, color=TEXT_GRAY if line else TEXT_GRAY, font_name="Consolas")


def create_closing_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ACCENT_PURPLE)

    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(1.0),
                 "💬 tong — Summary", font_size=48, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(0.8),
                 "A production-grade, real-time messaging platform with a comprehensive feature set",
                 font_size=20, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

    # Stats
    stats = [
        ("5", "Backend Controllers", ACCENT_BLUE),
        ("5", "Route Modules", ACCENT_GREEN),
        ("4", "Socket Handlers", ACCENT_PURPLE),
        ("4", "Backend Services", ACCENT_AMBER),
        ("8", "PostgreSQL Tables", ACCENT_PINK),
        ("20+", "Socket.io Events", ACCENT_BLUE),
        ("40+", "REST Endpoints", ACCENT_GREEN),
        ("16", "React Components", ACCENT_RED),
    ]

    for i, (num, label, color) in enumerate(stats):
        x = Inches(0.8 + (i % 4) * 3.1)
        y = Inches(3.8 + (i // 4) * 1.5)
        card = add_shape(slide, x, y, Inches(2.7), Inches(1.2), BG_CARD, color)
        add_text_box(slide, x + Inches(0.1), y + Inches(0.15), Inches(2.5), Inches(0.5),
                     num, font_size=32, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.1), y + Inches(0.7), Inches(2.5), Inches(0.4),
                     label, font_size=13, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.5),
                 "Phase 1: Complete Messaging System  •  React.js + Node.js + Socket.io  •  PostgreSQL + MongoDB + Redis",
                 font_size=14, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)


# ─── Main ────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    print("📊 Generating slides...")

    create_title_slide(prs)
    print("  ✅ Title slide")

    create_toc_slide(prs)
    print("  ✅ Table of Contents")

    create_overview_slide(prs)
    print("  ✅ Project Overview & Roadmap")

    create_architecture_slide(prs)
    print("  ✅ System Architecture")

    create_tech_stack_slide(prs)
    print("  ✅ Tech Stack")

    create_db_schema_slide(prs)
    print("  ✅ Database Schema")

    create_auth_workflow_slide(prs)
    print("  ✅ Authentication Workflow")

    create_socket_flow_slide(prs)
    print("  ✅ Socket.io Workflow")

    create_chat_workflow_slide(prs)
    print("  ✅ Chat Workflow")

    create_group_chat_slide(prs)
    print("  ✅ Group Chat")

    create_media_slide(prs)
    print("  ✅ Media Upload")

    create_message_lifecycle_slide(prs)
    print("  ✅ Message Lifecycle")

    create_notification_slide(prs)
    print("  ✅ Notifications & Presence")

    create_frontend_arch_slide(prs)
    print("  ✅ Frontend Architecture")

    create_api_summary_slide(prs)
    print("  ✅ API Summary")

    create_deployment_slide(prs)
    print("  ✅ Deployment — Server Architecture")

    create_deployment_flow_slide(prs)
    print("  ✅ Deployment — Deploy Flow & Services")

    create_file_structure_slide(prs)
    print("  ✅ File Structure")

    create_closing_slide(prs)
    print("  ✅ Summary/Closing slide")

    output_path = os.path.join(os.path.dirname(__file__), "Tong_Project_Documentation.pptx")
    prs.save(output_path)
    print(f"\n🎉 Documentation saved to: {output_path}")
    print(f"   Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
